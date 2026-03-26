"""Breadcrumb navigation bar component."""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from design_constants import (
    SECTIONS, ACCENT_TEAL, WHITE, TEXT_MID, PRIMARY_DARK,
    SLIDE_WIDTH, BREADCRUMB_TOP, BREADCRUMB_HEIGHT, FONT_FAMILY,
    LIGHT_GRAY,
)


def add_breadcrumb(slide, active_section_index):
    """Add a breadcrumb navigation bar below the title bar.

    Args:
        slide: The slide to add the breadcrumb to.
        active_section_index: Index of the currently active section (0-based).
    """
    n = len(SECTIONS)
    bar_left = Inches(0)
    bar_top = BREADCRUMB_TOP
    bar_width = SLIDE_WIDTH
    bar_height = BREADCRUMB_HEIGHT

    # Background strip
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, bar_left, bar_top, bar_width, bar_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    # Thin bottom border line
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, bar_left, bar_top + bar_height - Emu(12700),
        bar_width, Emu(12700),
    )
    border.fill.solid()
    border.fill.fore_color.rgb = LIGHT_GRAY
    border.line.fill.background()

    # Calculate positions
    total_padding = Inches(1.2)  # left + right padding
    usable = int(SLIDE_WIDTH) - int(total_padding)
    separator_width = Inches(0.25)
    total_sep = int(separator_width) * (n - 1)
    item_width = (usable - total_sep) // n
    start_x = int(total_padding) // 2

    for i, section in enumerate(SECTIONS):
        x = start_x + i * (item_width + int(separator_width))
        is_active = (i == active_section_index)

        if is_active:
            # Teal pill background
            pill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(x + item_width // 8),
                bar_top + Emu(int(bar_height) // 6),
                Emu(item_width * 3 // 4),
                Emu(int(bar_height) * 2 // 3),
            )
            pill.fill.solid()
            pill.fill.fore_color.rgb = ACCENT_TEAL
            pill.line.fill.background()
            # Smaller corner radius
            pill.adjustments[0] = 0.25

        # Section label
        txBox = slide.shapes.add_textbox(
            Emu(x), bar_top, Emu(item_width), bar_height,
        )
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = section
        run.font.size = Pt(9)
        run.font.name = FONT_FAMILY
        run.font.bold = is_active
        run.font.color.rgb = WHITE if is_active else TEXT_MID
        # Vertical centering
        tf.auto_size = None
        p.space_before = Pt(0)
        p.space_after = Pt(0)

        # Arrow separator (except after last)
        if i < n - 1:
            sep_x = x + item_width
            sep_box = slide.shapes.add_textbox(
                Emu(sep_x), bar_top, separator_width, bar_height,
            )
            sep_tf = sep_box.text_frame
            sep_tf.word_wrap = False
            sep_p = sep_tf.paragraphs[0]
            sep_p.alignment = PP_ALIGN.CENTER
            sep_run = sep_p.add_run()
            sep_run.text = ">"
            sep_run.font.size = Pt(8)
            sep_run.font.name = FONT_FAMILY
            sep_run.font.color.rgb = TEXT_MID
            sep_run.font.bold = False
