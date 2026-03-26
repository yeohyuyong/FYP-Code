"""Reusable slide templates (A-G) for the DSO DIIM Presentation."""
import os
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from design_constants import (
    PRIMARY_DARK, PRIMARY_BLUE, ACCENT_TEAL, ACCENT_GOLD, LIGHT_GRAY, WHITE,
    RED_ALERT, GREEN_SUCCESS, TEXT_DARK, TEXT_MID,
    SLIDE_WIDTH, SLIDE_HEIGHT, TITLE_BAR_HEIGHT, BREADCRUMB_HEIGHT,
    CONTENT_TOP, CONTENT_LEFT, CONTENT_WIDTH, CONTENT_HEIGHT,
    FONT_FAMILY, MARGIN_LEFT,
)
from breadcrumb import add_breadcrumb


# ── Helpers ──────────────────────────────────────────────────────────────────

def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_bar(slide, title_text):
    """Add a dark title bar at the top of a content slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        SLIDE_WIDTH, TITLE_BAR_HEIGHT,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_DARK
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(0.1), SLIDE_WIDTH - Inches(1.2), TITLE_BAR_HEIGHT,
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_FAMILY


def _add_text_box(slide, left, top, width, height, text, font_size=16,
                  color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                  font_name=FONT_FAMILY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def _add_bullets(slide, left, top, width, height, items, font_size=15,
                 color=TEXT_DARK, spacing=Pt(6)):
    """Add a bulleted list to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.space_before = Pt(0)
        p.level = 0

        # Handle sub-items (tuples: (text, level))
        if isinstance(item, tuple):
            text, level = item
            p.level = level
        else:
            text = item

        run = p.add_run()
        run.text = f"  {text}"
        run.font.size = Pt(font_size - (1 * p.level))
        run.font.color.rgb = color
        run.font.name = FONT_FAMILY

        # Add bullet character
        bullet_run = p.runs[0]
        bullet_run.text = f"\u2022  {text}" if p.level == 0 else f"   \u25E6  {text}"

    return txBox


def _add_image_safe(slide, image_path, left, top, width=None, height=None):
    """Add image if it exists, otherwise add a placeholder box."""
    if os.path.exists(image_path):
        if width and height:
            slide.shapes.add_picture(image_path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(image_path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(image_path, left, top, height=height)
        else:
            slide.shapes.add_picture(image_path, left, top)
    else:
        # Placeholder
        w = width or Inches(5)
        h = height or Inches(3.5)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        box.line.color.rgb = TEXT_MID
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        fname = os.path.basename(image_path)
        run.text = f"[Insert image]\n{fname}"
        run.font.size = Pt(11)
        run.font.color.rgb = TEXT_MID
        run.font.name = FONT_FAMILY


def _add_callout_box(slide, left, top, width, height, text,
                     border_color=ACCENT_TEAL, bg_color=None):
    """Add a callout/highlight box."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if bg_color:
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xE0, 0xF7, 0xFA)
    box.line.color.rgb = border_color
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.color.rgb = TEXT_DARK
    run.font.name = FONT_FAMILY
    return box


# ── Template A: Title Slide ──────────────────────────────────────────────────

def create_title_slide(prs, title, subtitle, details, logo_path=None):
    """Template A: Dark background title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide, PRIMARY_DARK)

    # Decorative accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), Inches(2.0), Inches(9.333), Pt(4),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_TEAL
    line.line.fill.background()

    # Title
    _add_text_box(slide, Inches(1.5), Inches(2.3), Inches(10.333), Inches(1.5),
                  title, font_size=32, color=WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)

    # Subtitle
    _add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10.333), Inches(0.6),
                  subtitle, font_size=18, color=ACCENT_TEAL, bold=False,
                  alignment=PP_ALIGN.CENTER)

    # Details (student, supervisor, date)
    _add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10.333), Inches(1.5),
                  details, font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC),
                  alignment=PP_ALIGN.CENTER)

    # Bottom accent line
    line2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2), Inches(6.8), Inches(9.333), Pt(2),
    )
    line2.fill.solid()
    line2.fill.fore_color.rgb = ACCENT_GOLD
    line2.line.fill.background()

    # Logo if provided
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.3),
                                 height=Inches(0.8))

    return slide


# ── Template B: Section Divider ──────────────────────────────────────────────

def create_section_divider(prs, section_number, section_title, summary,
                           active_section):
    """Template B: Section divider with number + breadcrumb."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PRIMARY_DARK)

    # Section number (gold)
    _add_text_box(slide, Inches(1), Inches(1.8), Inches(11.333), Inches(1.2),
                  f"{section_number:02d}", font_size=60, color=ACCENT_GOLD,
                  bold=True, alignment=PP_ALIGN.CENTER)

    # Section title
    _add_text_box(slide, Inches(1), Inches(3.2), Inches(11.333), Inches(1.0),
                  section_title, font_size=34, color=WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)

    # Summary line
    _add_text_box(slide, Inches(2), Inches(4.4), Inches(9.333), Inches(0.8),
                  summary, font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC),
                  alignment=PP_ALIGN.CENTER)

    # Breadcrumb at bottom
    add_breadcrumb(slide, active_section)

    return slide


# ── Template C: Content Slide ────────────────────────────────────────────────

def create_content_slide(prs, title, active_section, content_func):
    """Template C: Standard content slide.

    content_func(slide, left, top, width, height) populates the content area.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_GRAY)

    _add_title_bar(slide, title)
    add_breadcrumb(slide, active_section)

    # Content area
    content_func(slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)

    return slide


# ── Template D: Two-Column Comparison ────────────────────────────────────────

def create_two_column(prs, title, active_section, left_header, right_header,
                      left_func, right_func, bottom_func=None):
    """Template D: Two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_GRAY)

    _add_title_bar(slide, title)
    add_breadcrumb(slide, active_section)

    col_width = Inches(5.6)
    gap = Inches(0.5)
    left_x = CONTENT_LEFT
    right_x = Emu(int(CONTENT_LEFT) + int(col_width) + int(gap))

    # Column headers
    for x, header, color in [(left_x, left_header, PRIMARY_BLUE),
                              (right_x, right_header, ACCENT_TEAL)]:
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, CONTENT_TOP,
            col_width, Inches(0.45),
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        tf = hdr.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT_FAMILY

    col_top = Emu(int(CONTENT_TOP) + int(Inches(0.55)))
    col_height = Inches(4.5) if bottom_func else Inches(5.2)

    left_func(slide, left_x, col_top, col_width, col_height)
    right_func(slide, right_x, col_top, col_width, col_height)

    if bottom_func:
        bottom_top = Emu(int(col_top) + int(col_height) + int(Inches(0.15)))
        bottom_func(slide, CONTENT_LEFT, bottom_top,
                     CONTENT_WIDTH, Inches(1.0))

    return slide


# ── Template E: Key Finding ──────────────────────────────────────────────────

def create_key_finding(prs, title, active_section, stat_text, stat_label,
                       body_text, callout_text=None):
    """Template E: Large stat number + callout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    _add_title_bar(slide, title)
    add_breadcrumb(slide, active_section)

    # Large stat
    _add_text_box(slide, Inches(1), CONTENT_TOP, Inches(11.333), Inches(1.5),
                  stat_text, font_size=52, color=ACCENT_GOLD, bold=True,
                  alignment=PP_ALIGN.CENTER)

    # Stat label
    _add_text_box(slide, Inches(2), Emu(int(CONTENT_TOP) + int(Inches(1.6))),
                  Inches(9.333), Inches(0.5),
                  stat_label, font_size=18, color=TEXT_MID,
                  alignment=PP_ALIGN.CENTER)

    # Body text
    _add_text_box(slide, Inches(1.5), Emu(int(CONTENT_TOP) + int(Inches(2.4))),
                  Inches(10.333), Inches(2.0),
                  body_text, font_size=16, color=TEXT_DARK,
                  alignment=PP_ALIGN.CENTER)

    # Callout box
    if callout_text:
        _add_callout_box(
            slide, Inches(2), Emu(int(CONTENT_TOP) + int(Inches(4.6))),
            Inches(9.333), Inches(1.0), callout_text,
        )

    return slide


# ── Template F: Chart + Annotation ───────────────────────────────────────────

def create_chart_annotation(prs, title, active_section, image_path,
                            annotations, image_width=None):
    """Template F: Chart on left + annotation bullets on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_GRAY)

    _add_title_bar(slide, title)
    add_breadcrumb(slide, active_section)

    img_w = image_width or Inches(7.5)
    img_h = Inches(4.8)
    _add_image_safe(slide, image_path, CONTENT_LEFT, CONTENT_TOP, img_w, img_h)

    # Annotation panel on right
    ann_left = Emu(int(CONTENT_LEFT) + int(img_w) + int(Inches(0.3)))
    ann_width = Emu(int(CONTENT_WIDTH) - int(img_w) - int(Inches(0.3)))
    _add_bullets(slide, ann_left, CONTENT_TOP, ann_width, img_h,
                 annotations, font_size=13, color=TEXT_DARK)

    return slide


# ── Template G: Excel Screenshot ─────────────────────────────────────────────

def create_excel_screenshot(prs, title, active_section, image_path,
                            caption, annotations=None):
    """Template G: Excel screenshot with caption and optional annotations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_GRAY)

    _add_title_bar(slide, title)
    add_breadcrumb(slide, active_section)

    # Screenshot image (centered, with shadow effect via border)
    img_w = Inches(10)
    img_h = Inches(4.5)
    img_left = Emu((int(SLIDE_WIDTH) - int(img_w)) // 2)
    _add_image_safe(slide, image_path, img_left, CONTENT_TOP, img_w, img_h)

    # Caption below image
    cap_top = Emu(int(CONTENT_TOP) + int(img_h) + int(Inches(0.15)))
    _add_text_box(slide, Inches(1.5), cap_top, Inches(10.333), Inches(0.8),
                  caption, font_size=13, color=TEXT_MID,
                  alignment=PP_ALIGN.CENTER)

    # Annotations as bullets below caption if provided
    if annotations:
        ann_top = Emu(int(cap_top) + int(Inches(0.5)))
        _add_bullets(slide, Inches(2), ann_top, Inches(9.333), Inches(1.5),
                     annotations, font_size=13)

    return slide


# ── Utility: Add table ───────────────────────────────────────────────────────

def add_styled_table(slide, left, top, width, height, headers, rows,
                     header_color=PRIMARY_BLUE):
    """Add a styled table to a slide."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths evenly
    col_width = int(width) // n_cols
    for i in range(n_cols):
        table.columns[i].width = Emu(col_width)

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = TEXT_DARK
                p.font.name = FONT_FAMILY
                p.alignment = PP_ALIGN.CENTER

    return table_shape
