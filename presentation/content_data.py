"""All slide content data for the DSO DIIM Presentation (38 slides)."""
import os
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from design_constants import (
    PRIMARY_DARK, PRIMARY_BLUE, ACCENT_TEAL, ACCENT_GOLD, LIGHT_GRAY, WHITE,
    TEXT_DARK, TEXT_MID, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH,
    CONTENT_HEIGHT, FONT_FAMILY, SLIDE_WIDTH, RED_ALERT, GREEN_SUCCESS,
)

# Base path for images
BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG = {
    'inop_covid': os.path.join(BASE, 'covid', 'Inoperability_Evolution.png'),
    'loss_covid': os.path.join(BASE, 'covid', 'Economic_Loss_Evolution.png'),
    'inop_manpower': os.path.join(BASE, 'manpower-disruption', 'Inoperability_Evolution.png'),
    'loss_manpower': os.path.join(BASE, 'manpower-disruption', 'Economic_Loss_Evolution.png'),
    'covid_close': os.path.join(BASE, 'simulations', 'results', 'covid_k5_close_rates.png'),
    'covid_ratio': os.path.join(BASE, 'simulations', 'results', 'covid_k5_ratio_distributions.png'),
    'manpower_close': os.path.join(BASE, 'simulations', 'results', 'manpower_k5_close_rates.png'),
    'best_at_k': os.path.join(BASE, 'simulations', 'results', 'enhanced_best_at_k.png'),
    'scatter_covid': os.path.join(BASE, 'simulations', 'results', 'scatter_covid_19_k5.png'),
    'scatter_manpower': os.path.join(BASE, 'simulations', 'results', 'scatter_manpower_k5.png'),
    'timing': os.path.join(BASE, 'simulations', 'results', 'intervention_timing_plot.png'),
    'pca_robust': os.path.join(BASE, 'simulations', 'results', 'pca_robustness_plot.png'),
    'heatmap': os.path.join(BASE, 'simulations', 'results', 'sensitivity_intervention_heatmap.png'),
    'boundary': os.path.join(BASE, 'simulations', 'results', 'decision_boundary_curve.png'),
    'enhanced_simplified': os.path.join(BASE, 'simulations', 'results', 'enhanced_simplified_methods.png'),
    'logo': os.path.join(BASE, 'fyp report', 'Title', 'logo.png'),
    'covid_dr_close': os.path.join(BASE, 'simulations', 'results', 'covid_dr_close_rates.png'),
    'manpower_dr_close': os.path.join(BASE, 'simulations', 'results', 'manpower_dr_close_rates.png'),
    'covid_dr_ratio': os.path.join(BASE, 'simulations', 'results', 'covid_dr_ratio_distributions.png'),
    'manpower_dr_ratio': os.path.join(BASE, 'simulations', 'results', 'manpower_dr_ratio_distributions.png'),
}


def build_all_slides(prs):
    """Build all 38 slides."""
    from slide_templates import (
        create_title_slide, create_section_divider, create_content_slide,
        create_two_column, create_key_finding, create_chart_annotation,
        create_excel_screenshot, add_styled_table,
        _add_text_box, _add_bullets, _add_image_safe, _add_callout_box,
    )

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION 1: INTRODUCTION (Slides 1-4)
    # ═══════════════════════════════════════════════════════════════════════

    # Slide 1: Title
    create_title_slide(
        prs,
        title="Economic Resilience Assessment Using the\nDynamic Inoperability Input-Output Model (DIIM)",
        subtitle="Final Year Project  \u2014  Presented to DSO",
        details="Yeoh Yu Yong\nSupervisor: Assoc Prof Santos\nNanyang Technological University\n2025",
        logo_path=IMG['logo'],
    )

    # Slide 2: The Problem
    def slide2_content(slide, left, top, width, height):
        _add_text_box(slide, left, top, width, Inches(0.6),
                      "Singapore\u2019s economy is interconnected \u2014 a shock to one sector cascades to all others.",
                      font_size=20, color=PRIMARY_DARK, bold=True)

        bullets = [
            "COVID-19 shutdowns disrupted construction, which deprived building-material manufacturers of demand",
            "Transport restrictions cascaded into logistics and food supply chains",
            "Grey zone conflicts could trigger rapid foreign worker departures, disrupting labour-dependent sectors",
        ]
        _add_bullets(slide, left, Emu(int(top) + int(Inches(1.0))),
                     width, Inches(2.5), bullets, font_size=15)

        # Key question callout
        _add_callout_box(
            slide, Inches(1.5), Emu(int(top) + int(Inches(3.8))),
            Inches(10.333), Inches(1.2),
            "\U0001F50D  Key Question: How do we quickly identify which sectors to prioritise for intervention during a crisis?",
            border_color=ACCENT_GOLD,
        )

    create_content_slide(prs, "The Problem", 0, slide2_content)

    # Slide 3: Research Objectives
    def slide3_content(slide, left, top, width, height):
        objectives = [
            ("1", "Apply the DIIM to two disruption scenarios:\nCOVID-19 pandemic  &  Manpower disruption (grey zone conflict)"),
            ("2", "Benchmark five simplified sector-ranking methods\nagainst the full DIIM simulation"),
            ("3", "Determine when simplified methods are sufficient\nvs when the full DIIM is needed"),
        ]
        for i, (num, text) in enumerate(objectives):
            y = Emu(int(top) + i * int(Inches(1.7)))
            # Number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left, y, Inches(0.6), Inches(0.6),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = ACCENT_TEAL
            circle.line.fill.background()
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = num
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = WHITE
            run.font.name = FONT_FAMILY

            # Text
            _add_text_box(slide, Emu(int(left) + int(Inches(0.9))), y,
                          Emu(int(width) - int(Inches(1.0))), Inches(1.2),
                          text, font_size=16, color=TEXT_DARK)

    create_content_slide(prs, "Research Objectives", 0, slide3_content)

    # Slide 4: Presentation Roadmap
    def slide4_content(slide, left, top, width, height):
        sections = [
            ("01", "DIIM Model", "Model construction\n& components"),
            ("02", "DIIM Results", "Inoperability &\neconomic loss"),
            ("03", "Simplified\nMethods", "5 alternative\nranking approaches"),
            ("04", "Monte Carlo", "2,000 trial\nrobustness test"),
            ("05", "Results", "Close rates &\nperformance"),
            ("06", "Recommend-\nations", "Decision\nframework"),
        ]
        box_w = Inches(1.7)
        box_h = Inches(2.2)
        gap = Inches(0.25)
        start_x = Emu((int(SLIDE_WIDTH) - len(sections) * int(box_w) - (len(sections) - 1) * int(gap)) // 2)
        y = Emu(int(top) + int(Inches(0.8)))

        for i, (num, title, desc) in enumerate(sections):
            x = Emu(int(start_x) + i * (int(box_w) + int(gap)))

            # Box
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = WHITE
            box.line.color.rgb = PRIMARY_BLUE
            box.line.width = Pt(1.5)

            # Number
            _add_text_box(slide, x, Emu(int(y) + int(Inches(0.15))), box_w, Inches(0.4),
                          num, font_size=22, color=ACCENT_GOLD, bold=True,
                          alignment=PP_ALIGN.CENTER)
            # Title
            _add_text_box(slide, x, Emu(int(y) + int(Inches(0.55))), box_w, Inches(0.7),
                          title, font_size=13, color=PRIMARY_DARK, bold=True,
                          alignment=PP_ALIGN.CENTER)
            # Desc
            _add_text_box(slide, x, Emu(int(y) + int(Inches(1.3))), box_w, Inches(0.7),
                          desc, font_size=10, color=TEXT_MID,
                          alignment=PP_ALIGN.CENTER)

            # Arrow between boxes
            if i < len(sections) - 1:
                arr_x = Emu(int(x) + int(box_w))
                _add_text_box(slide, arr_x, Emu(int(y) + int(Inches(0.8))),
                              gap, Inches(0.4), "\u25B6", font_size=14,
                              color=TEXT_MID, alignment=PP_ALIGN.CENTER)

    create_content_slide(prs, "Presentation Roadmap", 0, slide4_content)

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION 2: DIIM MODEL CONSTRUCTION (Slides 5-14)
    # ═══════════════════════════════════════════════════════════════════════

    # Slide 5: Section Divider
    create_section_divider(prs, 1, "The DIIM Model",
                           "Construction of the Dynamic Inoperability Input-Output Model", 1)

    # Slide 6: What is Inoperability?
    def slide6_content(slide, left, top, width, height):
        _add_text_box(slide, left, top, width, Inches(0.6),
                      "Inoperability measures the normalised production loss of each sector",
                      font_size=18, color=PRIMARY_DARK, bold=True)

        bullets = [
            "Derived from the Leontief input-output model",
            "Measures how disruptions affect the functioning of each sector over time",
            "Values range from 0 (fully operational) to 1 (complete shutdown)",
        ]
        _add_bullets(slide, left, Emu(int(top) + int(Inches(0.9))),
                     Inches(5.5), Inches(2.0), bullets, font_size=15)

        # Diagram: planned vs degraded
        for i, (label, color, val) in enumerate([
            ("Planned Production", PRIMARY_BLUE, "100%"),
            ("Degraded Production", RED_ALERT, "60%"),
        ]):
            bx = Emu(int(left) + int(Inches(7)))
            by = Emu(int(top) + int(Inches(0.5)) + i * int(Inches(2.0)))
            bar_full_w = Inches(4)
            bar_h = Inches(0.6)
            ratio = 1.0 if i == 0 else 0.6

            # Bar background
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, bar_full_w, bar_h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = LIGHT_GRAY
            bg.line.color.rgb = TEXT_MID
            bg.line.width = Pt(0.5)

            # Filled portion
            filled = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, bx, by,
                Emu(int(bar_full_w * ratio)), bar_h,
            )
            filled.fill.solid()
            filled.fill.fore_color.rgb = color
            filled.line.fill.background()

            _add_text_box(slide, bx, Emu(int(by) - int(Inches(0.3))),
                          bar_full_w, Inches(0.3),
                          f"{label} ({val})", font_size=11, color=TEXT_DARK)

        # Inoperability formula
        _add_callout_box(
            slide, Emu(int(left) + int(Inches(6.5))), Emu(int(top) + int(Inches(3.5))),
            Inches(5), Inches(1.0),
            "Inoperability = (Planned - Degraded) / Planned\n= (100% - 60%) / 100% = 0.40",
            border_color=ACCENT_TEAL,
        )

    create_content_slide(prs, "What is Inoperability?", 1, slide6_content)

    # Slide 7: From IIM to DIIM
    def slide7_content(slide, left, top, width, height):
        _add_text_box(slide, left, top, Inches(11), Inches(0.5),
                      "Static IIM  \u2192  Dynamic DIIM", font_size=20,
                      color=PRIMARY_DARK, bold=True)

        # IIM box
        iim_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Emu(int(top) + int(Inches(0.8))),
            Inches(5.2), Inches(2.0),
        )
        iim_box.fill.solid()
        iim_box.fill.fore_color.rgb = WHITE
        iim_box.line.color.rgb = PRIMARY_BLUE
        iim_box.line.width = Pt(2)

        _add_text_box(slide, Emu(int(left) + int(Inches(0.2))),
                      Emu(int(top) + int(Inches(0.9))),
                      Inches(4.8), Inches(0.4),
                      "IIM (Static)", font_size=16, color=PRIMARY_BLUE, bold=True)
        _add_text_box(slide, Emu(int(left) + int(Inches(0.2))),
                      Emu(int(top) + int(Inches(1.4))),
                      Inches(4.8), Inches(0.4),
                      "q = A*q + c*", font_size=22, color=TEXT_DARK, bold=True,
                      alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, Emu(int(left) + int(Inches(0.2))),
                      Emu(int(top) + int(Inches(2.0))),
                      Inches(4.8), Inches(0.4),
                      "Equilibrium state only \u2014 no time dynamics",
                      font_size=12, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

        # Arrow
        _add_text_box(slide, Emu(int(left) + int(Inches(5.4))),
                      Emu(int(top) + int(Inches(1.3))),
                      Inches(1), Inches(0.5),
                      "\u27A1", font_size=28, color=ACCENT_GOLD,
                      alignment=PP_ALIGN.CENTER)

        # DIIM box
        diim_x = Emu(int(left) + int(Inches(6.2)))
        diim_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, diim_x, Emu(int(top) + int(Inches(0.8))),
            Inches(5.5), Inches(2.0),
        )
        diim_box.fill.solid()
        diim_box.fill.fore_color.rgb = WHITE
        diim_box.line.color.rgb = ACCENT_TEAL
        diim_box.line.width = Pt(2)

        _add_text_box(slide, Emu(int(diim_x) + int(Inches(0.2))),
                      Emu(int(top) + int(Inches(0.9))),
                      Inches(5.1), Inches(0.4),
                      "DIIM (Dynamic)", font_size=16, color=ACCENT_TEAL, bold=True)
        _add_text_box(slide, Emu(int(diim_x) + int(Inches(0.2))),
                      Emu(int(top) + int(Inches(1.4))),
                      Inches(5.1), Inches(0.4),
                      "q(t+1) = q(t) + K[A*q(t) + c* \u2212 q(t)]",
                      font_size=20, color=TEXT_DARK, bold=True,
                      alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, Emu(int(diim_x) + int(Inches(0.2))),
                      Emu(int(top) + int(Inches(2.0))),
                      Inches(5.1), Inches(0.4),
                      "Adds time variable t and recovery matrix K",
                      font_size=12, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

        # Key note
        bullets = [
            "K is a diagonal matrix \u2014 recovery depends only on each sector\u2019s own capacity",
            "Larger diagonal elements = stronger recovery capacity",
            "t is discrete time; q(t) is the inoperability vector at time t",
        ]
        _add_bullets(slide, left, Emu(int(top) + int(Inches(3.3))),
                     width, Inches(2.0), bullets, font_size=14)

    create_content_slide(prs, "From IIM to DIIM", 1, slide7_content)

    # Slide 8: DIIM Formula Overview
    def slide8_content(slide, left, top, width, height):
        # Main equation
        _add_text_box(slide, Inches(1), top, Inches(11.333), Inches(0.8),
                      "q(t+1) = q(t) + K [ A* q(t) + c* \u2212 q(t) ]",
                      font_size=28, color=PRIMARY_DARK, bold=True,
                      alignment=PP_ALIGN.CENTER)

        # Component labels
        components = [
            ("q(t)", "Inoperability\nvector at time t", ACCENT_TEAL),
            ("K", "Recovery capacity\n(diagonal matrix)", PRIMARY_BLUE),
            ("A*", "Interdependency\nmatrix (from IO table)", ACCENT_GOLD),
            ("c*", "Demand\nperturbation vector", RED_ALERT),
            ("q(0)", "Initial shock\n(starting condition)", GREEN_SUCCESS),
        ]

        box_w = Inches(2.1)
        box_h = Inches(1.6)
        start_x = Emu((int(SLIDE_WIDTH) - len(components) * int(box_w) - (len(components) - 1) * int(Inches(0.15))) // 2)
        comp_y = Emu(int(top) + int(Inches(1.3)))

        for i, (symbol, desc, color) in enumerate(components):
            x = Emu(int(start_x) + i * (int(box_w) + int(Inches(0.15))))

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, comp_y, box_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = WHITE
            box.line.color.rgb = color
            box.line.width = Pt(2)

            _add_text_box(slide, x, Emu(int(comp_y) + int(Inches(0.15))),
                          box_w, Inches(0.5), symbol, font_size=22,
                          color=color, bold=True, alignment=PP_ALIGN.CENTER)
            _add_text_box(slide, x, Emu(int(comp_y) + int(Inches(0.7))),
                          box_w, Inches(0.7), desc, font_size=11,
                          color=TEXT_DARK, alignment=PP_ALIGN.CENTER)

        # "We will show how each is obtained"
        _add_callout_box(
            slide, Inches(2.5), Emu(int(top) + int(Inches(3.5))),
            Inches(8.333), Inches(0.8),
            "\u27A1  We will now show how each component is obtained from real data.",
            border_color=ACCENT_GOLD,
        )

        # Two-phase note
        _add_text_box(slide, left, Emu(int(top) + int(Inches(4.6))),
                      width, Inches(0.4),
                      "Two phases:  Lockdown (c* \u2260 0)  \u2192  Recovery (c* = 0)",
                      font_size=14, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

    create_content_slide(prs, "DIIM Formula: Components", 1, slide8_content)

    # Slide 9: Obtaining A*
    create_excel_screenshot(
        prs, "Obtaining A* \u2014 Interdependency Matrix", 1,
        "",  # placeholder - user inserts IO table screenshot
        "A* = x\u207b\u00b9 A x  where A is the technical coefficient matrix from the IO table and x is total output",
        annotations=[
            "15\u00d715 Input-Output Table from Singapore\u2019s Department of Statistics",
            "A = raw IO flows / total output of purchasing sector",
            "A* normalises A by sector size for inoperability space",
            "2019 IO table used for COVID-19; 2022 IO table for Manpower scenario",
        ],
    )

    # Slide 10: Obtaining K
    create_excel_screenshot(
        prs, "Obtaining K \u2014 Recovery Capacity Matrix", 1,
        "",  # placeholder
        "K is a diagonal matrix where each K_ii represents sector i\u2019s recovery speed",
        annotations=[
            "K is diagonal: sector recovery depends only on its own production capacity",
            "Larger K_ii = faster recovery from disruption",
            "Estimated from historical recovery patterns (Santos et al., 2009)",
            "Terminal condition: inoperability at end = 1% of initial inoperability",
        ],
    )

    # Slide 11: Obtaining q(0) - COVID-19
    create_excel_screenshot(
        prs, "Obtaining q(0) \u2014 COVID-19 Initial Inoperability", 1,
        "",  # placeholder - user inserts unemployment data screenshot
        "q(0) estimated from sector-level unemployment rates during COVID-19 lockdown",
        annotations=[
            "Source: Ministry of Manpower unemployment data (19 years)",
            "Linear regression + 4-year moving average to estimate \u201cnormal\u201d 2020 rate",
            "q(0) = (actual rate \u2212 expected rate) / expected rate",
            "Mapped from MOM sectors to IOT 15-sector classification",
        ],
    )

    # Slide 12: Obtaining c*
    create_excel_screenshot(
        prs, "Obtaining c* \u2014 Demand Perturbation", 1,
        "",  # placeholder - user inserts GDP screenshot
        "c* captures the demand-side shock during the lockdown period; c* = 0 during recovery",
        annotations=[
            "Estimated from GDP decline data (SingStat Table Builder)",
            "c* represents reduced final demand during disruption",
            "Active only during lockdown phase; set to 0 during recovery",
            "For Manpower scenario: c* = 0 (pure supply-side shock, no demand perturbation)",
        ],
    )

    # Slide 13: Obtaining q(0) - Manpower
    create_excel_screenshot(
        prs, "Obtaining q(0) \u2014 Manpower Disruption", 1,
        "",  # placeholder
        "q(0) based on foreign worker dependency ratios by sector (2022 data)",
        annotations=[
            "Source: SG Foreign Manpower Dependency 2022",
            "Sectors with higher foreign worker share \u2192 higher initial inoperability",
            "2022 IO table aggregated to 15 sectors for consistency",
            "Key difference from COVID: c* = 0 (pure supply shock from worker departures)",
        ],
    )

    # Slide 14: COVID Timeline & Assumptions
    def slide14_content(slide, left, top, width, height):
        _add_text_box(slide, left, top, width, Inches(0.5),
                      "Model Parameters & Assumptions",
                      font_size=18, color=PRIMARY_DARK, bold=True)

        # Timeline visual
        tl_top = Emu(int(top) + int(Inches(1.0)))
        tl_left = Inches(1.5)
        tl_width = Inches(10)

        # Timeline bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, tl_left, tl_top, tl_width, Inches(0.15),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY_BLUE
        bar.line.fill.background()

        # Lockdown segment (red)
        lock_w = Emu(int(tl_width) * 55 // (55 + 974))
        lock = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, tl_left, tl_top, lock_w, Inches(0.15),
        )
        lock.fill.solid()
        lock.fill.fore_color.rgb = RED_ALERT
        lock.line.fill.background()

        # Labels
        _add_text_box(slide, tl_left, Emu(int(tl_top) - int(Inches(0.35))),
                      lock_w, Inches(0.3), "Lockdown\n7 Apr \u2013 1 Jun 2020 (55 days)",
                      font_size=10, color=RED_ALERT, bold=True,
                      alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, Emu(int(tl_left) + int(lock_w)),
                      Emu(int(tl_top) - int(Inches(0.35))),
                      Emu(int(tl_width) - int(lock_w)), Inches(0.3),
                      "Recovery\n2 Jun 2020 \u2013 1 Feb 2023 (974 days)",
                      font_size=10, color=PRIMARY_BLUE, bold=True,
                      alignment=PP_ALIGN.CENTER)

        # Assumptions
        bullets = [
            "No further shocks after initial lockdown period",
            "Terminal inoperability = 1% of initial inoperability (Santos et al., 2009)",
            "Production distributed evenly over the year",
            "Lockdown phase: c* \u2260 0 (demand shock active)",
            "Recovery phase: c* = 0 (demand returns to normal)",
        ]
        _add_bullets(slide, left, Emu(int(tl_top) + int(Inches(0.8))),
                     Inches(5.5), Inches(3.0), bullets, font_size=14)

        # Manpower comparison box
        _add_callout_box(
            slide, Emu(int(left) + int(Inches(6.5))),
            Emu(int(tl_top) + int(Inches(0.8))),
            Inches(5), Inches(2.5),
            "Manpower Scenario Differences:\n"
            "\u2022  2022 IO table (vs 2019)\n"
            "\u2022  c* = 0 throughout (pure supply shock)\n"
            "\u2022  q(0) from foreign worker dependency\n"
            "\u2022  Same recovery framework applies",
            border_color=ACCENT_TEAL,
        )

    create_content_slide(prs, "Timeline & Assumptions", 1, slide14_content)

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION 3: DIIM RESULTS (Slides 15-19)
    # ═══════════════════════════════════════════════════════════════════════

    # Slide 15: Section Divider
    create_section_divider(prs, 2, "DIIM Simulation Results",
                           "Inoperability propagation and economic loss across sectors", 2)

    # Slide 16: Inoperability Evolution
    create_chart_annotation(
        prs, "Inoperability Evolution \u2014 COVID-19", 2,
        IMG['inop_covid'],
        [
            "Each line = one sector\u2019s inoperability over time",
            "Sharp rise during lockdown (c* \u2260 0)",
            "Gradual recovery after lockdown ends (c* = 0)",
            "Some sectors recover faster (higher K_ii)",
            "DIIM captures the full cascade dynamics",
        ],
    )

    # Slide 17: Economic Loss Calculation
    def slide17_content(slide, left, top, width, height):
        _add_text_box(slide, left, top, width, Inches(0.5),
                      "Economic loss per sector per time step:",
                      font_size=18, color=PRIMARY_DARK, bold=True)

        # Formula
        _add_text_box(slide, Inches(2), Emu(int(top) + int(Inches(0.8))),
                      Inches(9.333), Inches(0.8),
                      "Loss_i(t) = x_i \u00d7 q_i(t) / 365",
                      font_size=28, color=ACCENT_TEAL, bold=True,
                      alignment=PP_ALIGN.CENTER)

        components = [
            ("x_i", "Daily output of sector i under normal production"),
            ("q_i(t)", "Inoperability of sector i at time t"),
            ("365", "Days in a year (evenly distributed production)"),
        ]
        for i, (sym, desc) in enumerate(components):
            y = Emu(int(top) + int(Inches(2.0)) + i * int(Inches(0.6)))
            _add_text_box(slide, Inches(2.5), y, Inches(1.5), Inches(0.5),
                          sym, font_size=16, color=ACCENT_GOLD, bold=True)
            _add_text_box(slide, Inches(4.2), y, Inches(6), Inches(0.5),
                          f"=  {desc}", font_size=15, color=TEXT_DARK)

        _add_text_box(slide, left, Emu(int(top) + int(Inches(4.2))),
                      width, Inches(0.5),
                      "Total economic loss = \u2211 Loss_i(t)  summed over all sectors and time steps",
                      font_size=15, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

    create_content_slide(prs, "Economic Loss Calculation", 2, slide17_content)

    # Slide 18: COVID vs Manpower Results (two-column)
    def left_18(slide, left, top, width, height):
        _add_image_safe(slide, IMG['loss_covid'], left, top, width, height)

    def right_18(slide, left, top, width, height):
        _add_image_safe(slide, IMG['loss_manpower'], left, top, width, height)

    def bottom_18(slide, left, top, width, height):
        from slide_templates import add_styled_table
        add_styled_table(
            slide, left, top, width, height,
            headers=["Parameter", "COVID-19", "Manpower"],
            rows=[
                ["IO Table", "2019", "2022"],
                ["Shock Type", "Supply + Demand", "Supply Only"],
                ["c*", "Non-zero (lockdown)", "Zero"],
                ["q(0) Source", "Unemployment rates", "Foreign worker dependency"],
            ],
        )

    create_two_column(prs, "COVID-19 vs Manpower: Economic Loss", 2,
                      "COVID-19 Scenario", "Manpower Scenario",
                      left_18, right_18, bottom_18)

    # Slide 19: Key Question
    create_key_finding(
        prs, "The Challenge", 2,
        "?",
        "Can we identify critical sectors WITHOUT running the full DIIM?",
        "The DIIM requires crisis-time data (q(0), c*) that may not be available "
        "during an emergency. Simplified methods use only the IO table structure, "
        "which is always available.",
        "If simplified methods match DIIM\u2019s recommendations \u2192 "
        "faster, more practical decision-making in crisis",
    )

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION 4: SIMPLIFIED METHODS (Slides 20-24)
    # ═══════════════════════════════════════════════════════════════════════

    # Slide 20: Section Divider
    create_section_divider(prs, 3, "Simplified Sector-Ranking Methods",
                           "Five methods that use only the IO table to rank sectors", 3)

    # Slide 21: Five Simplified Methods
    def slide21_content(slide, left, top, width, height):
        from slide_templates import add_styled_table
        add_styled_table(
            slide, left, top, Inches(11.5), Inches(3.8),
            headers=["Method", "Formula", "Intuition"],
            rows=[
                ["Total Output", "rank by x_i", "Biggest sectors matter most"],
                ["PCA \u00d7 x_i", "PC distance \u00d7 x_i", "Structural influence + size"],
                ["PageRank \u00d7 x_i", "PageRank centrality \u00d7 x_i", "Network importance + size"],
                ["BL \u00d7 x_i", "Backward linkage \u00d7 x_i", "Demand-pull influence + size"],
                ["FL \u00d7 x_i", "Forward linkage \u00d7 x_i", "Supply-push influence + size"],
            ],
            header_color=ACCENT_TEAL,
        )

        _add_callout_box(
            slide, Inches(2), Emu(int(top) + int(Inches(4.2))),
            Inches(9.333), Inches(0.9),
            "\U0001F4A1  Key insight: All methods are weighted by total output (x_i) to reflect "
            "economic size. This means economic size is always a factor in the ranking.",
            border_color=ACCENT_GOLD,
        )

    create_content_slide(prs, "Five Simplified Methods", 3, slide21_content)

    # Slide 22: How Methods Are Compared
    def slide22_content(slide, left, top, width, height):
        steps = [
            ("1", "Compute\nSimplified\nRanking", "Static, no\nsimulation needed", PRIMARY_BLUE),
            ("2", "Select\nTop-k\nSectors", "Choose k most\nimportant sectors", ACCENT_TEAL),
            ("3", "Run DIIM\nwith\nIntervention", "Reduce inoperability\nfor selected sectors", ACCENT_GOLD),
            ("4", "Compare\nLoss\nReduction", "vs DIIM\u2019s own\ntop-k selection", GREEN_SUCCESS),
        ]

        box_w = Inches(2.4)
        box_h = Inches(2.5)
        gap = Inches(0.4)
        start_x = Emu((int(SLIDE_WIDTH) - len(steps) * int(box_w) - (len(steps) - 1) * int(gap)) // 2)

        for i, (num, title, desc, color) in enumerate(steps):
            x = Emu(int(start_x) + i * (int(box_w) + int(gap)))
            y = top

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = WHITE
            box.line.color.rgb = color
            box.line.width = Pt(2)

            _add_text_box(slide, x, Emu(int(y) + int(Inches(0.15))),
                          box_w, Inches(0.5), f"Step {num}", font_size=14,
                          color=color, bold=True, alignment=PP_ALIGN.CENTER)
            _add_text_box(slide, x, Emu(int(y) + int(Inches(0.6))),
                          box_w, Inches(0.9), title, font_size=15,
                          color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
            _add_text_box(slide, x, Emu(int(y) + int(Inches(1.5))),
                          box_w, Inches(0.8), desc, font_size=11,
                          color=TEXT_MID, alignment=PP_ALIGN.CENTER)

            if i < len(steps) - 1:
                _add_text_box(slide, Emu(int(x) + int(box_w)),
                              Emu(int(y) + int(Inches(0.9))),
                              gap, Inches(0.5), "\u27A1", font_size=22,
                              color=TEXT_MID, alignment=PP_ALIGN.CENTER)

        # Metric
        _add_callout_box(
            slide, Inches(2), Emu(int(top) + int(Inches(3.2))),
            Inches(9.333), Inches(1.6),
            "Performance Metric:  R_m = (Simplified reduction) / (DIIM reduction)\n\n"
            "R_m = 1.0 means the simplified method matches DIIM perfectly\n"
            "R_m < 0.8 means the simplified method misses >20% of DIIM\u2019s benefit",
            border_color=ACCENT_TEAL,
        )

    create_content_slide(prs, "Evaluation Pipeline", 3, slide22_content)

    # Slide 23: Method Rankings Preview
    create_excel_screenshot(
        prs, "Method Rankings \u2014 Side by Side", 3,
        "",  # placeholder
        "Each column shows one method\u2019s sector ranking. Arrows highlight where rankings agree or disagree.",
        annotations=[
            "Rankings are fixed (computed once from IO table)",
            "These static rankings are then tested against 2,000 random disruption profiles",
            "Key observation: most methods agree on the top 2\u20133 sectors",
        ],
    )

    # Slide 24: Loss Reduction Algorithm
    def slide24_content(slide, left, top, width, height):
        _add_text_box(slide, left, top, width, Inches(0.5),
                      "How intervention reduces economic loss:",
                      font_size=18, color=PRIMARY_DARK, bold=True)

        steps = [
            "1. For each selected sector i, reduce its inoperability: q_i(0) \u2192 q_i(0) \u00d7 0.5",
            "2. Re-run DIIM with the modified q(0) vector",
            "3. Compute total economic loss under intervention",
            "4. Loss reduction = baseline loss \u2212 intervention loss",
        ]
        _add_bullets(slide, left, Emu(int(top) + int(Inches(0.8))),
                     Inches(6), Inches(3.0), steps, font_size=15)

        _add_callout_box(
            slide, Emu(int(left) + int(Inches(6.5))), Emu(int(top) + int(Inches(0.5))),
            Inches(5), Inches(3.0),
            "The simplified ranking is fixed \u2014 it doesn\u2019t change with q(0).\n\n"
            "But DIIM\u2019s optimal ranking can change with each disruption profile.\n\n"
            "Question: How often does the simplified ranking still pick\n"
            "the \u201cright\u201d sectors?",
            border_color=ACCENT_TEAL,
        )

        _add_text_box(slide, left, Emu(int(top) + int(Inches(4.3))),
                      width, Inches(0.5),
                      "\u27A1  Answer: Monte Carlo simulation with 2,000 random disruption profiles",
                      font_size=16, color=ACCENT_GOLD, bold=True,
                      alignment=PP_ALIGN.CENTER)

    create_content_slide(prs, "Loss Reduction Algorithm", 3, slide24_content)

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION 5: MONTE CARLO FRAMEWORK (Slides 25-27)
    # ═══════════════════════════════════════════════════════════════════════

    # Slide 25: Section Divider
    create_section_divider(prs, 4, "Monte Carlo Evaluation",
                           "Testing robustness with 2,000 random disruption profiles", 4)

    # Slide 26: Monte Carlo Design
    def slide26_content(slide, left, top, width, height):
        steps = [
            ("1", "Start with real q(0) baseline", PRIMARY_BLUE),
            ("2", "Perturb: q(0)_rand = q(0) \u00d7 exp(N(0, 0.5))", ACCENT_TEAL),
            ("3", "Run DIIM baseline + DIIM with simplified top-k", ACCENT_GOLD),
            ("4", "Compute R_m = simplified / DIIM reduction", RED_ALERT),
            ("5", "Repeat 2,000 times", PRIMARY_BLUE),
        ]

        for i, (num, text, color) in enumerate(steps):
            y = Emu(int(top) + i * int(Inches(0.9)))
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left, y, Inches(0.5), Inches(0.5),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.fill.background()
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = num
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = WHITE
            run.font.name = FONT_FAMILY

            _add_text_box(slide, Emu(int(left) + int(Inches(0.7))), y,
                          Inches(5.5), Inches(0.5),
                          text, font_size=15, color=TEXT_DARK)

        # Metrics box
        metrics_x = Emu(int(left) + int(Inches(7)))
        _add_callout_box(
            slide, metrics_x, top, Inches(4.5), Inches(3.5),
            "Performance Metrics:\n\n"
            "\u2022  Mean Ratio (R\u0304_m):\n"
            "   Average R_m across all 2,000 trials\n\n"
            "\u2022  80% Close Rate:\n"
            "   Fraction of trials where R_m \u2265 0.80\n"
            "   (simplified captures \u226580% of DIIM benefit)",
            border_color=ACCENT_GOLD,
        )

    create_content_slide(prs, "Monte Carlo Framework", 4, slide26_content)

    # Slide 27: Three Performance Regimes
    def slide27_content(slide, left, top, width, height):
        regimes = [
            ("Small k (k=3)", "Methods diverge\nsignificantly",
             "Sector selection\nmatters most", RED_ALERT, "R\u0304_m: 0.65\u20130.87"),
            ("Medium k (k=5\u20137)", "Sweet spot for\ndifferentiation",
             "Best methods\nidentifiable", ACCENT_GOLD, "Close rate: 55\u201390%"),
            ("Large k (k\u226510)", "All methods\nconverge",
             "Simplified\nsufficient", GREEN_SUCCESS, "Close rate: ~100%"),
        ]

        box_w = Inches(3.5)
        box_h = Inches(3.5)
        gap = Inches(0.4)
        start_x = Emu((int(SLIDE_WIDTH) - len(regimes) * int(box_w) - (len(regimes) - 1) * int(gap)) // 2)

        for i, (title, desc, implication, color, stat) in enumerate(regimes):
            x = Emu(int(start_x) + i * (int(box_w) + int(gap)))

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = WHITE
            box.line.color.rgb = color
            box.line.width = Pt(3)

            _add_text_box(slide, x, Emu(int(top) + int(Inches(0.2))),
                          box_w, Inches(0.5), title, font_size=18,
                          color=color, bold=True, alignment=PP_ALIGN.CENTER)
            _add_text_box(slide, x, Emu(int(top) + int(Inches(0.8))),
                          box_w, Inches(0.8), desc, font_size=14,
                          color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
            _add_text_box(slide, x, Emu(int(top) + int(Inches(1.7))),
                          box_w, Inches(0.5), stat, font_size=13,
                          color=color, bold=True, alignment=PP_ALIGN.CENTER)
            _add_text_box(slide, x, Emu(int(top) + int(Inches(2.4))),
                          box_w, Inches(0.8), implication, font_size=13,
                          color=TEXT_MID, alignment=PP_ALIGN.CENTER)

        _add_text_box(slide, left, Emu(int(top) + int(Inches(3.8))),
                      width, Inches(0.5),
                      "This is the central finding: the value of method selection depends on the intervention scale k",
                      font_size=14, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

    create_content_slide(prs, "Three Performance Regimes", 4, slide27_content)

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION 6: RESULTS (Slides 28-34)
    # ═══════════════════════════════════════════════════════════════════════

    # Slide 28: Section Divider
    create_section_divider(prs, 5, "Results",
                           "Performance comparison across scenarios and intervention scales", 5)

    # Slide 29: COVID-19 Close Rates
    create_chart_annotation(
        prs, "COVID-19: Close Rates by k", 5,
        IMG['covid_close'],
        [
            "Close rate = fraction of trials where R_m \u2265 0.80",
            "At k=5: best methods achieve 89.8% close rate",
            "At k\u226510: all methods converge to ~100%",
            "PCA\u00d7x_i shows slight advantage at small k",
        ],
    )

    # Slide 30: COVID-19 Distributions
    create_chart_annotation(
        prs, "COVID-19: Ratio Distributions", 5,
        IMG['covid_ratio'],
        [
            "Distribution of R_m across 2,000 trials",
            "Tighter distributions = more consistent performance",
            "Most distributions cluster near 1.0 for k\u22657",
        ],
    )

    # Slide 31: Manpower Close Rates
    create_chart_annotation(
        prs, "Manpower: Close Rates by k", 5,
        IMG['manpower_close'],
        [
            "At k=5: FL\u00d7x_i and Total Output lead at 55.4%",
            "Supply-only shocks are harder to approximate",
            "Faster convergence: ~97% close rate by k=7",
            "c*=0 makes dynamics more predictable",
        ],
    )

    # Slide 32: Head-to-Head
    def left_32(slide, left, top, width, height):
        _add_image_safe(slide, IMG['best_at_k'], left, top, width, height)

    def right_32(slide, left, top, width, height):
        bullets = [
            "Economic size (Total Output) dominates both scenarios when weighted by x_i",
            "At large k: specific ranking is immaterial",
            "COVID-19 shows more method differentiation (due to c* \u2260 0)",
            "Manpower converges faster (pure supply dynamics more predictable)",
        ]
        _add_bullets(slide, left, top, width, height, bullets, font_size=13)

    create_two_column(prs, "Head-to-Head: COVID vs Manpower", 5,
                      "Best Method at Each k", "Key Observations",
                      left_32, right_32)

    # Slide 33: Scatter
    create_chart_annotation(
        prs, "Scatter: Simplified vs DIIM Loss Reduction", 5,
        IMG['scatter_covid'],
        [
            "Each point = one Monte Carlo trial",
            "45\u00b0 line = perfect match with DIIM",
            "Points above line: simplified outperforms",
            "Points below: DIIM would have been better",
            "Clustering near the line = good agreement",
        ],
    )

    # Slide 34: Robustness Evidence
    def slide34_content(slide, left, top, width, height):
        _add_text_box(slide, left, top, width, Inches(0.4),
                      "Results hold across timing, parameter, and profile variations",
                      font_size=16, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

        # 2x2 grid of images
        img_w = Inches(5.5)
        img_h = Inches(2.3)
        gap = Inches(0.3)
        images = [
            (IMG['timing'], "Intervention Timing"),
            (IMG['pca_robust'], "PCA Robustness"),
            (IMG['heatmap'], "Sensitivity Heatmap"),
            (IMG['boundary'], "Decision Boundary"),
        ]
        for i, (path, caption) in enumerate(images):
            row = i // 2
            col = i % 2
            x = Emu(int(left) + col * (int(img_w) + int(gap)))
            y = Emu(int(top) + int(Inches(0.6)) + row * (int(img_h) + int(Inches(0.4))))
            _add_image_safe(slide, path, x, y, img_w, img_h)
            _add_text_box(slide, x, Emu(int(y) + int(img_h)),
                          img_w, Inches(0.3), caption, font_size=10,
                          color=TEXT_MID, alignment=PP_ALIGN.CENTER)

    create_content_slide(prs, "Robustness Evidence", 5, slide34_content)

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION 7: RECOMMENDATIONS (Slides 35-38)
    # ═══════════════════════════════════════════════════════════════════════

    # Slide 35: Section Divider
    create_section_divider(prs, 6, "Recommendations",
                           "When to use simplified methods vs full DIIM", 6)

    # Slide 36: Decision Framework
    def slide36_content(slide, left, top, width, height):
        _add_text_box(slide, left, top, width, Inches(0.6),
                      "How many sectors to intervene in?",
                      font_size=22, color=PRIMARY_DARK, bold=True,
                      alignment=PP_ALIGN.CENTER)

        # Three decision boxes
        decisions = [
            ("k \u2265 10", "Use simplified method (any)",
             "All methods converge \u2014\nspecific ranking is immaterial",
             GREEN_SUCCESS, "\u2713"),
            ("k = 5\u20137", "Use simplified, validate with DIIM",
             "Methods differ \u2014\nvalidate for critical decisions",
             ACCENT_GOLD, "\u26A0"),
            ("k < 5", "Run full DIIM",
             "Simplified methods are insufficient \u2014\n"
             "sector selection matters too much",
             RED_ALERT, "\u2717"),
        ]

        box_w = Inches(3.5)
        box_h = Inches(2.5)
        gap = Inches(0.3)
        start_x = Emu((int(SLIDE_WIDTH) - len(decisions) * int(box_w) - (len(decisions) - 1) * int(gap)) // 2)
        dec_y = Emu(int(top) + int(Inches(1.0)))

        for i, (k_range, action, reason, color, icon) in enumerate(decisions):
            x = Emu(int(start_x) + i * (int(box_w) + int(gap)))

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, dec_y, box_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = WHITE
            box.line.color.rgb = color
            box.line.width = Pt(3)

            _add_text_box(slide, x, Emu(int(dec_y) + int(Inches(0.15))),
                          box_w, Inches(0.5), f"{icon}  {k_range}",
                          font_size=22, color=color, bold=True,
                          alignment=PP_ALIGN.CENTER)
            _add_text_box(slide, x, Emu(int(dec_y) + int(Inches(0.7))),
                          box_w, Inches(0.5), action, font_size=14,
                          color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
            _add_text_box(slide, x, Emu(int(dec_y) + int(Inches(1.3))),
                          box_w, Inches(1.0), reason, font_size=12,
                          color=TEXT_MID, alignment=PP_ALIGN.CENTER)

        # Bottom summary
        _add_callout_box(
            slide, Inches(2), Emu(int(dec_y) + int(Inches(2.8))),
            Inches(9.333), Inches(0.8),
            "Peacetime / large-scale planning \u2192 Simplified methods sufficient    |    "
            "Targeted / crisis response \u2192 Full DIIM required",
            border_color=ACCENT_GOLD,
        )

    create_content_slide(prs, "Decision Framework", 6, slide36_content)

    # Slide 37: Summary
    def slide37_content(slide, left, top, width, height):
        _add_text_box(slide, left, top, width, Inches(0.5),
                      "Key Contributions", font_size=20, color=PRIMARY_DARK, bold=True)

        contributions = [
            ("1", "Extended DIIM to two distinct disruption scenarios",
             "COVID-19 (supply + demand) and Manpower (pure supply)", ACCENT_TEAL),
            ("2", "Demonstrated three performance regimes",
             "Small k: methods diverge | Medium k: sweet spot | Large k: convergence", ACCENT_GOLD),
            ("3", "Provided actionable decision framework",
             "When to use simplified vs full DIIM based on intervention scale k", PRIMARY_BLUE),
        ]

        for i, (num, title, desc, color) in enumerate(contributions):
            y = Emu(int(top) + int(Inches(0.8)) + i * int(Inches(1.5)))
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left, y, Inches(0.6), Inches(0.6),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.fill.background()
            tf = circle.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = num
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = WHITE
            run.font.name = FONT_FAMILY

            _add_text_box(slide, Emu(int(left) + int(Inches(0.9))), y,
                          Emu(int(width) - int(Inches(1.0))), Inches(0.4),
                          title, font_size=16, color=TEXT_DARK, bold=True)
            _add_text_box(slide, Emu(int(left) + int(Inches(0.9))),
                          Emu(int(y) + int(Inches(0.5))),
                          Emu(int(width) - int(Inches(1.0))), Inches(0.5),
                          desc, font_size=13, color=TEXT_MID)

    create_content_slide(prs, "Summary", 6, slide37_content)

    # Slide 38: Future Work & Thank You
    def slide38_content(slide, left, top, width, height):
        _add_text_box(slide, left, top, Inches(5.5), Inches(0.5),
                      "Future Directions", font_size=18, color=PRIMARY_DARK, bold=True)

        futures = [
            "Broader disruption scenarios (supply chain, cyber, climate)",
            "Dynamic simplified methods that adapt to observed shocks",
            "Higher sector granularity (beyond 15 sectors)",
            "Cross-country comparison with other IO tables",
        ]
        _add_bullets(slide, left, Emu(int(top) + int(Inches(0.6))),
                     Inches(5.5), Inches(3.0), futures, font_size=14)

        # Thank you
        _add_text_box(slide, Emu(int(left) + int(Inches(6.5))), top,
                      Inches(5), Inches(1.0),
                      "Thank You", font_size=36, color=ACCENT_GOLD, bold=True,
                      alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, Emu(int(left) + int(Inches(6.5))),
                      Emu(int(top) + int(Inches(1.5))),
                      Inches(5), Inches(1.5),
                      "Questions & Discussion",
                      font_size=20, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

        # Contact
        _add_text_box(slide, Emu(int(left) + int(Inches(6.5))),
                      Emu(int(top) + int(Inches(3.0))),
                      Inches(5), Inches(1.0),
                      "Yeoh Yu Yong\nNanyang Technological University",
                      font_size=14, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

    create_content_slide(prs, "Future Work & Thank You", 6, slide38_content)
